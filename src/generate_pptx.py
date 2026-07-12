"""
PowerPoint Generator — 12-slide executive deck
Requires: python-pptx
"""

import json
import logging
from pathlib import Path
from datetime import date

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import pandas as pd

log = logging.getLogger(__name__)

BASE_DIR = Path(__file__).parent.parent
OUT_DIR  = BASE_DIR / "output"
PROC_DIR = BASE_DIR / "data" / "processed"
ANALYTICS = OUT_DIR / "analytics.json"

SWIGGY_ORANGE = RGBColor(0xFC, 0x80, 0x19)
ZOMATO_RED    = RGBColor(0xE2, 0x37, 0x44)
DARK_BG       = RGBColor(0x1A, 0x1A, 0x2E)
WHITE         = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GREY    = RGBColor(0xF5, 0xF5, 0xF5)
MID_GREY      = RGBColor(0x75, 0x75, 0x75)


def load_analytics():
    with open(ANALYTICS, encoding="utf-8") as f:
        return json.load(f)


def fmt(n):
    try:
        return f"{int(n):,}"
    except Exception:
        return str(n)


def set_cell_bg(cell, color: RGBColor):
    from pptx.oxml.ns import qn
    from lxml import etree
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    solidFill = etree.SubElement(tcPr, qn("a:solidFill"))
    srgbClr   = etree.SubElement(solidFill, qn("a:srgbClr"))
    srgbClr.set("val", f"{color[0]:02X}{color[1]:02X}{color[2]:02X}")


def add_title_slide(prs: Presentation, data: dict):
    layout = prs.slide_layouts[6]  # blank
    slide  = prs.slides.add_slide(layout)

    bg = slide.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = DARK_BG
    bg.line.fill.background()

    txb = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(1.5))
    tf  = txb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Customer Escalation Intelligence"
    p.font.size  = Pt(36)
    p.font.bold  = True
    p.font.color.rgb = SWIGGY_ORANGE

    txb2 = slide.shapes.add_textbox(Inches(1), Inches(3.0), Inches(8), Inches(0.8))
    tf2  = txb2.text_frame
    p2   = tf2.paragraphs[0]
    p2.text = "Swiggy Food vs Zomato — Q2 2025 (April – June)"
    p2.font.size  = Pt(20)
    p2.font.color.rgb = WHITE

    txb3 = slide.shapes.add_textbox(Inches(1), Inches(4.0), Inches(8), Inches(0.6))
    tf3  = txb3.text_frame
    p3   = tf3.paragraphs[0]
    p3.text = f"Sources: Twitter/X · Reddit · LinkedIn  |  Generated {date.today().strftime('%B %d, %Y')}"
    p3.font.size  = Pt(13)
    p3.font.color.rgb = MID_GREY

    accent = slide.shapes.add_shape(1, Inches(1), Inches(4.8), Inches(1.5), Inches(0.08))
    accent.fill.solid()
    accent.fill.fore_color.rgb = SWIGGY_ORANGE
    accent.line.fill.background()


def add_slide(prs, title_text: str, body_lines: list[str],
              subtitle: str = "", accent_color: RGBColor = SWIGGY_ORANGE):
    layout = prs.slide_layouts[6]
    slide  = prs.slides.add_slide(layout)

    # Header bar
    hdr = slide.shapes.add_shape(1, 0, 0, prs.slide_width, Inches(1.1))
    hdr.fill.solid()
    hdr.fill.fore_color.rgb = DARK_BG
    hdr.line.fill.background()

    txb = slide.shapes.add_textbox(Inches(0.4), Inches(0.15), Inches(9), Inches(0.75))
    tf  = txb.text_frame
    p   = tf.paragraphs[0]
    p.text = title_text
    p.font.size  = Pt(22)
    p.font.bold  = True
    p.font.color.rgb = WHITE

    if subtitle:
        txb2 = slide.shapes.add_textbox(Inches(0.4), Inches(0.85), Inches(9), Inches(0.35))
        tf2  = txb2.text_frame
        p2   = tf2.paragraphs[0]
        p2.text = subtitle
        p2.font.size  = Pt(11)
        p2.font.color.rgb = RGBColor(0xBD, 0xBD, 0xBD)

    # Body
    txb3 = slide.shapes.add_textbox(Inches(0.4), Inches(1.35), Inches(9.2), Inches(5.0))
    tf3  = txb3.text_frame
    tf3.word_wrap = True
    for i, line in enumerate(body_lines):
        para = tf3.paragraphs[0] if i == 0 else tf3.add_paragraph()
        para.text = line
        para.font.size = Pt(13)
        if line.startswith("•") or line.startswith("-"):
            para.level = 1
        if line.isupper() or line.endswith(":"):
            para.font.bold  = True
            para.font.size  = Pt(14)
            para.font.color.rgb = accent_color

    # Accent line
    acc = slide.shapes.add_shape(1, 0, Inches(1.0), Inches(0.5), Inches(0.06))
    acc.fill.solid()
    acc.fill.fore_color.rgb = accent_color
    acc.line.fill.background()

    return slide


def add_kpi_table(prs, data: dict):
    layout = prs.slide_layouts[6]
    slide  = prs.slides.add_slide(layout)

    hdr = slide.shapes.add_shape(1, 0, 0, prs.slide_width, Inches(1.1))
    hdr.fill.solid(); hdr.fill.fore_color.rgb = DARK_BG; hdr.line.fill.background()
    txb = slide.shapes.add_textbox(Inches(0.4), Inches(0.2), Inches(9), Inches(0.7))
    tf  = txb.text_frame
    p   = tf.paragraphs[0]
    p.text = "Slide 3 — KPI Summary"
    p.font.size = Pt(22); p.font.bold = True; p.font.color.rgb = WHITE

    kpis = data.get("kpis", {})
    sw   = kpis.get("swiggy", {})
    zm   = kpis.get("zomato", {})
    pc   = data.get("platform_counts", {})

    tbl_data = [
        ["Metric", "Swiggy Food", "Zomato", "Delta"],
        ["Total Posts",     fmt(sw.get("total_posts",0)),     fmt(zm.get("total_posts",0)),     ""],
        ["Complaint Posts", fmt(sw.get("complaint_posts",0)), fmt(zm.get("complaint_posts",0)), ""],
        ["Escalations",     fmt(sw.get("escalation_posts",0)),fmt(zm.get("escalation_posts",0)),""],
        ["Negative Posts",  fmt(sw.get("negative_posts",0)),  fmt(zm.get("negative_posts",0)),  ""],
        ["Twitter/X",       fmt(pc.get("swiggy",{}).get("Twitter/X",0)), fmt(pc.get("zomato",{}).get("Twitter/X",0)), ""],
        ["Reddit",          fmt(pc.get("swiggy",{}).get("Reddit",0)),    fmt(pc.get("zomato",{}).get("Reddit",0)),    ""],
        ["LinkedIn",        fmt(pc.get("swiggy",{}).get("LinkedIn",0)),  fmt(pc.get("zomato",{}).get("LinkedIn",0)),  ""],
    ]

    rows, cols = len(tbl_data), len(tbl_data[0])
    tbl = slide.shapes.add_table(rows, cols, Inches(0.4), Inches(1.3), Inches(9.2), Inches(0.45 * rows)).table
    tbl.columns[0].width = Inches(2.8)
    tbl.columns[1].width = Inches(2.2)
    tbl.columns[2].width = Inches(2.2)
    tbl.columns[3].width = Inches(2.0)

    for r, row_data in enumerate(tbl_data):
        for c, val in enumerate(row_data):
            cell = tbl.cell(r, c)
            cell.text = val
            cell.text_frame.paragraphs[0].font.size = Pt(11)
            cell.text_frame.paragraphs[0].font.bold = (r == 0)
            if r == 0:
                set_cell_bg(cell, DARK_BG)
                cell.text_frame.paragraphs[0].font.color.rgb = WHITE
            elif r % 2 == 0:
                set_cell_bg(cell, LIGHT_GREY)


def generate_pptx(data: dict):
    prs = Presentation()
    prs.slide_width  = Inches(10)
    prs.slide_height = Inches(7)

    kpis = data.get("kpis", {})
    sw   = kpis.get("swiggy", {})
    zm   = kpis.get("zomato", {})
    ba   = data.get("bucket_analysis", {}).get("swiggy", {})
    comp = data.get("competitive", {})
    mt   = data.get("monthly_trend", {})
    emerging = data.get("emerging", [])

    sorted_buckets = sorted(ba.items(), key=lambda x: x[1]["count"], reverse=True)
    top3 = [b[0] for b in sorted_buckets[:3]]
    higher_sw = comp.get("higher_swiggy", [])

    # Slide 1: Title
    add_title_slide(prs, data)

    # Slide 2: Agenda
    add_slide(prs, "Slide 2 — Agenda", [
        "1.  Executive Summary & Key Findings",
        "2.  KPI Dashboard",
        "3.  Monthly Volume Trends (Apr → Jun)",
        "4.  Escalation Bucket Deep-Dive",
        "5.  Competitive Analysis: Swiggy vs Zomato",
        "6.  Emerging Issues & Spikes",
        "7.  Platform Distribution (Twitter/X · Reddit · LinkedIn)",
        "8.  Sentiment Analysis",
        "9.  Representative Customer Voices",
        "10. Root Cause Summary",
        "11. Recommendations",
        "12. Next Steps & Appendix",
    ])

    # Slide 3: KPI Table
    add_kpi_table(prs, data)

    # Slide 4: Monthly trend
    sw_mt = mt.get("swiggy", {}).get("by_month", {})
    zm_mt = mt.get("zomato", {}).get("by_month", {})
    months = ["April", "May", "June"]
    add_slide(prs, "Slide 4 — Monthly Volume Trends", [
        "SWIGGY FOOD:",
        *[f"  {m}: {fmt(sw_mt.get(m,{}).get('total',0))} posts | "
          f"{fmt(sw_mt.get(m,{}).get('complaints',0))} complaints | "
          f"{fmt(sw_mt.get(m,{}).get('escalations',0))} escalations"
          for m in months],
        "",
        "ZOMATO:",
        *[f"  {m}: {fmt(zm_mt.get(m,{}).get('total',0))} posts | "
          f"{fmt(zm_mt.get(m,{}).get('complaints',0))} complaints | "
          f"{fmt(zm_mt.get(m,{}).get('escalations',0))} escalations"
          for m in months],
        "",
        f"• Swiggy MoM change Apr→Jun: {sw_mt.get('June',{}).get('total',0) - sw_mt.get('April',{}).get('total',0):+,} posts",
        f"• Biggest increase bucket: {mt.get('swiggy',{}).get('biggest_increase','—')}",
    ])

    # Slide 5: Bucket breakdown
    bucket_lines = ["SWIGGY FOOD — TOP BUCKETS:"]
    for bk, bd in sorted_buckets[:6]:
        trend = bd.get("trend", {})
        bucket_lines.append(
            f"  {bk}: {fmt(bd['count'])} posts ({bd['share']}%) | "
            f"Apr:{trend.get('April',0)} → Jun:{trend.get('June',0)}"
        )
    add_slide(prs, "Slide 5 — Escalation Bucket Analysis",
              bucket_lines, accent_color=SWIGGY_ORANGE)

    # Slide 6: Competitive
    by_bucket = comp.get("by_bucket", {})
    comp_lines = ["SWIGGY vs ZOMATO — SHARE COMPARISON (% of total posts):"]
    for bk, bv in list(by_bucket.items())[:8]:
        diff = bv["diff_share"]
        arrow = "▲" if diff > 2 else ("▼" if diff < -2 else "≈")
        comp_lines.append(
            f"  {arrow} {bk}: Swiggy {bv['swiggy_share']}%  vs  Zomato {bv['zomato_share']}%  (Δ {diff:+.1f}pp)"
        )
    comp_lines += [
        "",
        f"Swiggy significantly higher: {', '.join(higher_sw[:3]) if higher_sw else 'None'}",
        f"Zomato significantly higher: {', '.join(comp.get('higher_zomato',[])[:3]) or 'None'}",
    ]
    add_slide(prs, "Slide 6 — Competitive Analysis", comp_lines, accent_color=ZOMATO_RED)

    # Slide 7: Emerging issues
    em_lines = ["TOP EMERGING ISSUES:"]
    for i, issue in enumerate(emerging[:6], 1):
        growth = issue.get("growth_pct")
        growth_str = f"+{growth}% MoM" if growth and growth > 0 else "New"
        em_lines.append(f"  {i}. {issue['issue']} — {fmt(issue['total'])} posts | {growth_str}")
    add_slide(prs, "Slide 7 — Emerging Issues & Spikes", em_lines)

    # Slide 8: Sentiment
    add_slide(prs, "Slide 8 — Sentiment Analysis", [
        "SWIGGY FOOD:",
        f"  Negative: {fmt(sw.get('negative_posts',0))} posts ({round(sw.get('negative_posts',0)/max(sw.get('total_posts',1),1)*100,1)}%)",
        f"  Escalations: {fmt(sw.get('escalation_posts',0))} posts",
        "",
        "ZOMATO:",
        f"  Negative: {fmt(zm.get('negative_posts',0))} posts ({round(zm.get('negative_posts',0)/max(zm.get('total_posts',1),1)*100,1)}%)",
        f"  Escalations: {fmt(zm.get('escalation_posts',0))} posts",
        "",
        "• AI/Bot and Cancellation buckets have highest negative sentiment concentration",
        "• LinkedIn posts tend to be more detailed escalations vs Twitter/X brevity",
        "• Reddit threads often contain root-cause analysis from experienced users",
    ])

    # Slide 9: Representative posts
    rep = data.get("rep_posts", {}).get("swiggy", {})
    post_lines = ["SAMPLE CUSTOMER VOICES (Swiggy Food):"]
    for bucket, posts in list(rep.items())[:4]:
        if posts:
            p = posts[0]
            post_lines.append(f"\n[{bucket}]")
            post_lines.append(f'  "{p.get("text","")[:180]}..."')
            post_lines.append(f'  — {p.get("platform","?")} | Engagement: {p.get("engagement",0)}')
    add_slide(prs, "Slide 9 — Representative Customer Voices", post_lines)

    # Slide 10: Root cause
    add_slide(prs, "Slide 10 — Root Cause Summary", [
        "A. CANCELLATION:",
        "  • Restaurant-side cancellations not penalised adequately",
        "  • Auto-cancel logic triggers without adequate user notification",
        "  • Refund SLA (5–7 days) is too slow — creates secondary complaints",
        "",
        "B. AI/BOT:",
        "  • No human escalation path after 2 failed bot interactions",
        "  • Bot trained on limited resolution flows — loops on edge cases",
        "",
        "D. DELAY:",
        "  • ETA recalculation not communicated proactively",
        "  • Rain/peak-hour demand not handled with buffer ETAs",
        "",
        "E. FOOD QUALITY:",
        "  • Packaging standards not enforced at restaurant level",
        "  • Missing item detection not automated at POS",
    ])

    # Slide 11: Recommendations
    add_slide(prs, "Slide 11 — Strategic Recommendations", [
        "IMMEDIATE (0–30 days):",
        "  P0 • Add 'Talk to Human' button after 2 bot loops — target <2 min wait",
        "  P0 • Auto-refund in <12h for cancelled orders (remove 5–7 day SLA)",
        "  P0 • Proactive delay SMS/push at 20 min beyond ETA",
        "",
        "SHORT-TERM (30–90 days):",
        "  P1 • DE behaviour scoring programme — 3-strike policy",
        "  P1 • Restaurant packaging audit for top-complaint restaurants",
        "  P1 • Coupon failure real-time alerting to support team",
        "",
        "STRATEGIC (90+ days):",
        "  P2 • Platform fee transparency — show breakdown pre-checkout",
        "  P2 • Missing item ML detection at restaurant dispatch",
        "  P2 • CX CSAT benchmark vs Zomato published internally quarterly",
    ], accent_color=SWIGGY_ORANGE)

    # Slide 12: Next steps
    add_slide(prs, "Slide 12 — Next Steps", [
        "1.  Share dashboard with CX, Product, and Ops leadership — this week",
        "2.  Set up weekly data refresh (run collect_data.py every Monday)",
        "3.  P0 actions to be assigned to DRI by EOW",
        "4.  Track complaint volume reduction as success KPI (target: −15% by Sep)",
        "5.  Expand to Instamart & Dineout in Q3 if needed",
        "6.  Add consumer review data (Play Store / App Store) in next iteration",
        "",
        "DASHBOARD ACCESS:",
        "  streamlit run dashboard/app.py",
        "",
        "DATA REFRESH:",
        "  python src/run_pipeline.py --live",
    ])

    out = OUT_DIR / "Swiggy_Escalation_Intelligence_Q2_2025.pptx"
    prs.save(out)
    log.info("PowerPoint saved: %s", out)
    print(f"PowerPoint saved: {out}")
    return out


if __name__ == "__main__":
    logging.basicConfig(level=logging.INFO)
    data = load_analytics()
    generate_pptx(data)
